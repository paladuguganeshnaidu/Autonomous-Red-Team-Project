from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
TEMPLATE = Path(r"c:\Users\ganes\Downloads\[EXT] Solution Challenge 2026 - Prototype PPT Template (2).pptx")
OUTPUT = ROOT / "deliverables" / "Autonomous_Red_Team_Agent_Deck.pptx"


SLIDE_TEXT = {
    1: {
        "title": "Autonomous Red Team Agent",
        "body": [
            "Conversational reconnaissance and reporting assistant for authorized security assessments.",
            "The agent plans scans, runs recon tools, analyzes results, and explains what it did in plain language.",
            "Current prototype supports CLI and web chat workflows with a final written report after each run.",
        ],
    },
    2: {
        "title": "Team Details",
        "body": [
            "Team name: Autonomous Red Team Project",
            "Team leader name: Ganes",
            "Problem Statement: Security testing tools are powerful but fragmented, technical, and hard to narrate for operators in real time.",
        ],
    },
    3: {
        "title": "Brief About The Solution",
        "body": [
            "This solution combines an orchestration loop, recon tooling, memory, and report generation into one operator-facing agent.",
            "A user can ask for a basic scan in natural language, watch the agent describe each command it runs, and receive a final report in the same conversation.",
            "The goal is to reduce manual context-switching while keeping scan activity explainable and auditable.",
        ],
    },
    4: {
        "title": "Opportunities",
        "body": [
            "How different is it: Instead of acting like a silent automation script, the agent communicates progress, reasoning, commands used, and final findings as a teammate.",
            "How it solves the problem: It turns multiple recon steps into one guided workflow from target input to evidence-backed report.",
            "USP: Deterministic basic recon, chat-based interaction, command history, persistent scan memory, and actionable reporting in one lightweight prototype.",
        ],
    },
    5: {
        "title": "List Of Features Offered By The Solution",
        "body": [
            "Natural-language scan requests from CLI or web chat",
            "Deterministic recon flow using subfinder, httpx, nmap, and ffuf",
            "LLM-assisted reasoning for intent classification and follow-up answers",
            "Live progress narration with exact commands executed",
            "Persistent session memory and reusable chat history",
            "Human-readable final report with findings, weaknesses, and next steps",
        ],
    },
    6: {
        "title": "Process Flow Diagram / Use Case",
        "body": [
            "Operator request -> target extraction -> scan plan selection",
            "Planner -> executor -> analyzer -> state persistence",
            "Tool outputs -> conversational updates -> report generation",
            "Operator follow-up questions -> memory lookup -> direct answer or LLM-assisted summary",
            "Primary use case: 'Make a basic scan on target.com and give me the report.'",
        ],
    },
    7: {
        "title": "Wireframes / Mock Diagrams",
        "body": [
            "Web UI: left sidebar for chat history, center chat stream for agent messages, input box for new operator requests.",
            "Live execution cards: show planning output, executing command, and finished summary for each scan step.",
            "Final interaction style: the agent replies like an analyst, not just a terminal window.",
        ],
    },
    8: {
        "title": "Architecture Diagram",
        "body": [
            "Frontend: FastAPI web app + WebSocket log streaming",
            "Orchestration core: main loop with planner, executor, analyzer, and state manager",
            "Recon tool layer: subfinder, nmap, httpx, ffuf wrappers",
            "Intelligence layer: local LLM/Ollama-compatible prompt handlers for planning, memory Q&A, and enrichment",
            "Outputs: chat transcript, session memory, logs, and final report artifacts",
        ],
    },
    9: {
        "title": "Technologies To Be Used In The Solution",
        "body": [
            "Python 3.13",
            "FastAPI + Uvicorn for web serving",
            "PowerShell / shell execution for tool orchestration",
            "subfinder, nmap, httpx, ffuf for recon",
            "Ollama-compatible LLM endpoint for reasoning",
            "JSON-based session memory, logs, and report artifacts",
        ],
    },
    10: {
        "title": "Estimated Implementation Cost",
        "body": [
            "Prototype cost is low when run locally on an existing machine with open-source tooling.",
            "Core expenses for production would be cloud hosting, model inference, storage, access control, and monitoring.",
            "Expected prototype budget: minimal software cost plus operator time for testing and validation.",
        ],
    },
    11: {
        "title": "Snapshots Of The MVP",
        "body": [
            "Live scan sample on toriiminds.com discovered services on ports 22, 80, 443, 3000, 3001, 3002, 4000, 4001, and 4002.",
            "The updated web agent now announces each scan step, keeps command history, and posts the final report back into the chat thread.",
            "Generated report includes recon summary, weaknesses, recommended next steps, and execution timeline.",
        ],
    },
    12: {
        "title": "Additional Details / Future Development",
        "body": [
            "Add richer evidence collection such as screenshots, nuclei integration, and endpoint screenshots for discovered panels.",
            "Support role-based access, multi-user collaboration, and per-target workspaces.",
            "For Solution Challenge packaging, the reasoning layer can be adapted to Gemini/Vertex AI and cloud-hosted execution controls.",
        ],
    },
    13: {
        "title": "Project Links",
        "body": [
            "GitHub Public Repository: Add your public repo link before submission",
            "Demo Video Link (3 Minutes): To be recorded",
            "MVP Link: Local prototype currently runs at http://127.0.0.1:8000",
            "Working Prototype Link: Publish the FastAPI app to cloud and attach URL here",
        ],
    },
}


def set_text_box(shape, title: str, body_lines: list[str]) -> None:
    text_frame = shape.text_frame
    text_frame.clear()
    text_frame.word_wrap = True

    p = text_frame.paragraphs[0]
    p.text = title
    p.font.bold = True
    p.font.size = Pt(24)
    p.alignment = PP_ALIGN.LEFT

    for line in body_lines:
        paragraph = text_frame.add_paragraph()
        paragraph.text = line
        paragraph.font.size = Pt(16)
        paragraph.level = 0
        paragraph.alignment = PP_ALIGN.LEFT


def first_text_shape(slide):
    for shape in slide.shapes:
        if hasattr(shape, "text_frame"):
            return shape
    return None


def clear_other_text_shapes(slide, keep_shape) -> None:
    for shape in slide.shapes:
        if shape is keep_shape:
            continue
        if hasattr(shape, "text_frame"):
            shape.text_frame.clear()


def add_text_slide(slide, title: str, body_lines: list[str]) -> None:
    textbox = slide.shapes.add_textbox(Inches(0.75), Inches(1.2), Inches(11.5), Inches(5.4))
    set_text_box(textbox, title, body_lines)


def build_presentation() -> Path:
    prs = Presentation(str(TEMPLATE))

    for slide_number, payload in SLIDE_TEXT.items():
        slide = prs.slides[slide_number - 1]
        shape = first_text_shape(slide)
        if shape is None:
            add_text_slide(slide, payload["title"], payload["body"])
        else:
            clear_other_text_shapes(slide, shape)
            set_text_box(shape, payload["title"], payload["body"])

    add_text_slide(
        prs.slides[13],
        "Demo Story",
        [
            "1. Operator asks for a basic scan on a target.",
            "2. Agent confirms the target and starts deterministic recon.",
            "3. UI streams the commands being run and summarizes each result.",
            "4. Agent produces a final report and answers follow-up questions like 'what did you run?'",
        ],
    )

    add_text_slide(
        prs.slides[14],
        "Thank You / Q&A",
        [
            "Autonomous Red Team Agent",
            "A transparent recon assistant that scans, explains, and reports in one workflow.",
            "Prepared for rapid prototype presentation.",
        ],
    )

    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUTPUT))
    return OUTPUT


if __name__ == "__main__":
    path = build_presentation()
    print(path)
